#!/usr/bin/env python3
from __future__ import annotations

import json
import os
from pathlib import Path
from typing import Iterable

from PIL import Image

# Use writable cache path for matplotlib in this environment.
ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentations"
ASSETS_DIR = PRESENTATION_DIR / "assets"
os.environ.setdefault("MPLCONFIGDIR", str(PRESENTATION_DIR / ".mplconfig"))

import matplotlib  # noqa: E402

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


REPORT_PATH = ROOT / "PROJECT_ANALYSIS_FULL.md"
CLEANING_SUMMARY_PATH = ROOT / "unzipped" / "cleaning" / "summary.json"
OOF_METRICS_PATH = ROOT / "outputs_top1_mps" / "all_model_oof_metrics.json"
REFINED_SUMMARY_PATH = ROOT / "outputs_top1_mps" / "analysis" / "refined_ensemble_summary.json"
PAIR_SUMMARY_PATH = ROOT / "outputs_top1_mps" / "analysis" / "pair_experts_integration_summary.json"
META_SUMMARY_PATH = ROOT / "outputs_top1_mps" / "analysis" / "meta_hgb_summary.json"
RISK_AUDIT_PATH = ROOT / "outputs_top1_mps" / "analysis" / "private_risk_audit.json"
COLOR_ABLATION_PATH = ROOT / "outputs_color_ablation_onefold_mps" / "one_fold_ablation_summary.json"

IMG_COVER = ROOT / "external_data_lab" / "reports" / "images" / "packeat_7482.jpg"
IMG_CLEANING = ROOT / "unzipped" / "cleaning" / "cleaning_preview_top_suspects.jpg"
IMG_AUG_DEBUG = ROOT / "outputs_top1_mps" / "analysis" / "aug_debug" / "orange_confusions_raw_vs_aug_sheet.png"

OUTPUT_PPTX = PRESENTATION_DIR / "DL_LAB1_MD_SHOWCASE_5_SLIDES.pptx"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# Palette
BG = RGBColor(11, 16, 38)
BG_ALT = RGBColor(16, 24, 58)
CARD = RGBColor(23, 34, 78)
CARD_2 = RGBColor(28, 45, 96)
ACCENT = RGBColor(27, 196, 176)
ACCENT_2 = RGBColor(255, 130, 65)
WHITE = RGBColor(242, 246, 255)
MUTED = RGBColor(181, 192, 225)


def _load_json(path: Path) -> dict:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def _style_paragraph(paragraph, size: int, color: RGBColor, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    paragraph.alignment = align
    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, base: RGBColor, add_blobs: bool = True) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = base
    bg.line.fill.background()
    if not add_blobs:
        return

    blobs = [
        (-0.8, -0.6, 3.6, 3.6, ACCENT, 0.88),
        (9.2, -0.9, 5.1, 3.8, ACCENT_2, 0.86),
        (10.8, 4.4, 4.2, 3.0, ACCENT, 0.90),
    ]
    for x, y, ww, hh, color, tr in blobs:
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(ww), Inches(hh))
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.fill.transparency = tr
        shp.line.fill.background()


def add_card(slide, x: float, y: float, w: float, h: float, color: RGBColor = CARD) -> None:
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.fill.transparency = 0.08
    card.line.fill.background()


def add_title(slide, text: str, x: float, y: float, w: float, h: float) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    _style_paragraph(tf.paragraphs[0], size=42, color=WHITE, bold=True)


def add_subtitle(slide, text: str, x: float, y: float, w: float, h: float, size: int = 18, color: RGBColor = MUTED) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    _style_paragraph(tf.paragraphs[0], size=size, color=color, bold=False)


def add_bullets(slide, lines: Iterable[str], x: float, y: float, w: float, h: float, size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {line}"
        p.level = 0
        _style_paragraph(p, size=size, color=WHITE)


def add_image_cover(slide, source_path: Path, x: float, y: float, w: float, h: float, temp_name: str) -> None:
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    target_ratio = w / h
    with Image.open(source_path) as img:
        src_w, src_h = img.size
        src_ratio = src_w / src_h
        if src_ratio > target_ratio:
            new_w = int(src_h * target_ratio)
            left = (src_w - new_w) // 2
            box = (left, 0, left + new_w, src_h)
        else:
            new_h = int(src_w / target_ratio)
            top = (src_h - new_h) // 2
            box = (0, top, src_w, top + new_h)
        cropped = img.crop(box).resize((int(w * 320), int(h * 320)), Image.Resampling.LANCZOS)
        out_path = ASSETS_DIR / temp_name
        cropped.save(out_path)
    slide.shapes.add_picture(str(out_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _plot_ensemble_progress(path: Path, refined: dict, pair: dict, meta: dict) -> None:
    labels = ["Base", "Weight\nSearch", "Bias", "Pair\nExperts", "Meta HGB"]
    values = [
        refined["base_metrics_from_existing_weights"]["acc"],
        refined["after_weight_search_metrics"]["acc"],
        refined["after_bias_metrics"]["acc"],
        pair["metrics_after_pair_experts"]["acc"],
        meta["meta_cv_metrics"]["acc"],
    ]
    colors = ["#5B74D6", "#7087E7", "#7FA3F3", "#2AC4B0", "#FF8241"]

    plt.figure(figsize=(8.4, 4.3), facecolor="#111933")
    ax = plt.gca()
    ax.set_facecolor("#111933")
    bars = ax.bar(labels, values, color=colors, width=0.62)
    ax.set_ylim(0.94, 0.98)
    ax.set_ylabel("Accuracy", color="#DDE6FF")
    ax.tick_params(colors="#DDE6FF")
    ax.grid(axis="y", color="#394675", alpha=0.5, linestyle="--")
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, val + 0.0005, f"{val:.4f}", ha="center", color="#F4F7FF", fontsize=9)
    plt.tight_layout()
    plt.savefig(path, dpi=220, facecolor="#111933")
    plt.close()


def _plot_risk(path: Path, risk: dict) -> None:
    labels = ["Meta-CV ACC", "Bootstrap p05", "Holdout min", "Group median", "Group min"]
    values = [
        risk["meta_cv"]["acc"],
        risk["bootstrap_oof"]["acc_p05"],
        risk["repeated_stratified_holdout"]["acc_min"],
        risk["group_holdout_by_plu"]["acc_median"],
        risk["group_holdout_by_plu"]["acc_min"],
    ]
    colors = ["#2AC4B0", "#46CFA4", "#82D58F", "#FFB160", "#FF5F5F"]

    plt.figure(figsize=(8.4, 4.2), facecolor="#121A38")
    ax = plt.gca()
    ax.set_facecolor("#121A38")
    y_pos = list(range(len(labels)))
    ax.barh(y_pos, values, color=colors, height=0.58)
    ax.set_yticks(y_pos)
    ax.set_yticklabels(labels, color="#DDE6FF")
    ax.set_xlim(0.35, 1.0)
    ax.tick_params(axis="x", colors="#DDE6FF")
    ax.grid(axis="x", color="#3A4A7E", alpha=0.5, linestyle="--")
    for y, v in zip(y_pos, values):
        ax.text(v + 0.01, y, f"{v:.3f}", va="center", color="#F4F7FF", fontsize=9)
    ax.invert_yaxis()
    plt.tight_layout()
    plt.savefig(path, dpi=220, facecolor="#121A38")
    plt.close()


def _plot_color_ablation(path: Path, ablation: dict) -> None:
    full = ablation["results"]["full"]["metrics_on_same_fold"]
    no_color = ablation["results"]["no_color"]["metrics_on_same_fold"]

    labels = ["ACC", "F1-macro"]
    full_vals = [full["acc"], full["f1_macro"]]
    no_color_vals = [no_color["acc"], no_color["f1_macro"]]
    x = [0, 1]
    width = 0.35

    plt.figure(figsize=(7.6, 4.2), facecolor="#131D3C")
    ax = plt.gca()
    ax.set_facecolor("#131D3C")
    ax.bar([xi - width / 2 for xi in x], full_vals, width=width, color="#6D80D8", label="full")
    ax.bar([xi + width / 2 for xi in x], no_color_vals, width=width, color="#2AC4B0", label="no_color")
    ax.set_xticks(x)
    ax.set_xticklabels(labels, color="#DDE6FF")
    ax.set_ylim(0.93, 0.965)
    ax.tick_params(axis="y", colors="#DDE6FF")
    ax.grid(axis="y", color="#3A4A7E", alpha=0.5, linestyle="--")
    for xi, fval, nval in zip(x, full_vals, no_color_vals):
        ax.text(xi - width / 2, fval + 0.0006, f"{fval:.4f}", ha="center", color="#F4F7FF", fontsize=9)
        ax.text(xi + width / 2, nval + 0.0006, f"{nval:.4f}", ha="center", color="#F4F7FF", fontsize=9)
    ax.legend(facecolor="#131D3C", edgecolor="#4E5B8F", labelcolor="#F4F7FF")
    plt.tight_layout()
    plt.savefig(path, dpi=220, facecolor="#131D3C")
    plt.close()


def build_presentation() -> Path:
    PRESENTATION_DIR.mkdir(parents=True, exist_ok=True)
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    cleaning = _load_json(CLEANING_SUMMARY_PATH)
    oof = _load_json(OOF_METRICS_PATH)
    refined = _load_json(REFINED_SUMMARY_PATH)
    pair = _load_json(PAIR_SUMMARY_PATH)
    meta = _load_json(META_SUMMARY_PATH)
    risk = _load_json(RISK_AUDIT_PATH)
    ablation = _load_json(COLOR_ABLATION_PATH)

    chart_ensemble = ASSETS_DIR / "chart_ensemble_progress.png"
    chart_risk = ASSETS_DIR / "chart_risk_profile.png"
    chart_ablation = ASSETS_DIR / "chart_color_ablation.png"
    _plot_ensemble_progress(chart_ensemble, refined, pair, meta)
    _plot_risk(chart_risk, risk)
    _plot_color_ablation(chart_ablation, ablation)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # Slide 1: Cover
    s1 = prs.slides.add_slide(blank)
    add_bg(s1, BG)
    add_card(s1, 0.7, 0.8, 7.0, 5.9, CARD)
    add_title(s1, "DL Lab 1\nОформление MD-отчёта\nв презентацию", 1.05, 1.1, 6.4, 2.2)
    add_subtitle(
        s1,
        "Источник: PROJECT_ANALYSIS_FULL.md\n"
        "Фокус: данные, пайплайн, метрики, риски и план следующего цикла",
        1.05,
        3.3,
        6.2,
        1.6,
        size=17,
    )
    add_subtitle(s1, "Дата отчёта: 2026-02-18/19", 1.05, 5.55, 4.0, 0.5, size=14, color=ACCENT)
    add_card(s1, 8.0, 0.8, 4.8, 5.9, CARD_2)
    add_image_cover(s1, IMG_COVER, 8.25, 1.05, 4.3, 5.2, "cover_image.jpg")

    # Slide 2: Data and cleaning
    s2 = prs.slides.add_slide(blank)
    add_bg(s2, BG_ALT)
    add_title(s2, "Данные и очистка train", 0.8, 0.45, 6.0, 0.9)
    add_card(s2, 0.8, 1.4, 5.3, 5.6, CARD)
    add_bullets(
        s2,
        [
            f"Всего train: {cleaning['total_rows']}",
            f"Drop: {cleaning['drop_count']} | Quarantine: {cleaning['quarantine_count']}",
            f"Strict: {cleaning['strict_count']} | Aggressive: {cleaning['aggressive_count']}",
            f"Основной drop-риск: too_small_min_side_lt_56 ({cleaning['drop_reason_counts']['too_small_min_side_lt_56']})",
            "Очистка применяется только к train, test не модифицируется",
        ],
        1.05,
        1.95,
        4.8,
        4.5,
        size=18,
    )
    add_card(s2, 6.35, 1.4, 6.1, 5.6, CARD_2)
    add_image_cover(s2, IMG_CLEANING, 6.6, 1.7, 5.6, 4.2, "cleaning_preview.jpg")
    add_subtitle(s2, "Визуальный артефакт очистки: top suspects", 6.72, 6.05, 5.5, 0.5, size=13, color=MUTED)

    # Slide 3: Stack and quality growth
    s3 = prs.slides.add_slide(blank)
    add_bg(s3, BG)
    add_title(s3, "Архитектура решения и прирост качества", 0.8, 0.45, 7.5, 0.9)
    add_card(s3, 0.8, 1.35, 4.25, 5.95, CARD)
    add_bullets(
        s3,
        [
            "3 backbone-модели, 5-fold CV (15 чекпоинтов)",
            f"ConvNeXt: {oof['convnext_small']['acc']:.4f} ACC",
            f"EffNetV2-S: {oof['effnetv2_s']['acc']:.4f} ACC",
            f"ResNet50: {oof['resnet50']['acc']:.4f} ACC",
            "Далее: weight search + class-bias + pair-experts + meta HGB stack",
        ],
        1.05,
        1.9,
        3.8,
        4.8,
        size=16,
    )
    add_card(s3, 5.25, 1.35, 7.2, 5.95, CARD_2)
    s3.shapes.add_picture(str(chart_ensemble), Inches(5.55), Inches(1.75), width=Inches(6.6), height=Inches(3.4))
    add_subtitle(
        s3,
        f"Meta stack: {meta['meta_cv_metrics']['acc']:.4f} ACC / {meta['meta_cv_metrics']['f1_macro']:.4f} F1",
        5.65,
        5.35,
        6.5,
        0.6,
        size=16,
        color=ACCENT,
    )
    add_subtitle(
        s3,
        f"Абсолютный прирост к refined-base: +{meta['delta_acc_abs']:.4f} ACC",
        5.65,
        5.9,
        6.5,
        0.5,
        size=14,
        color=MUTED,
    )

    # Slide 4: Risk and error profile
    s4 = prs.slides.add_slide(blank)
    add_bg(s4, BG_ALT)
    add_title(s4, "Риск private и профиль ошибок", 0.8, 0.45, 6.4, 0.9)
    add_card(s4, 0.8, 1.35, 7.0, 5.95, CARD)
    s4.shapes.add_picture(str(chart_risk), Inches(1.1), Inches(1.75), width=Inches(6.4), height=Inches(3.2))
    add_subtitle(
        s4,
        "Вывод: IID-оценки стабильные, но group-holdout по PLU показывает высокий\nриск domain shift (median 0.860, min 0.421).",
        1.1,
        5.1,
        6.3,
        1.2,
        size=14,
        color=WHITE,
    )
    add_card(s4, 8.05, 1.35, 4.4, 5.95, CARD_2)
    add_image_cover(s4, IMG_AUG_DEBUG, 8.35, 1.65, 3.8, 3.4, "aug_debug_cover.png")
    add_bullets(
        s4,
        [
            "Топ-путаницы: Мандарины↔Апельсин, Яблоки красные↔Томаты",
            "Hard-группы: Лук/197, Яблоки красные/807",
            "Это больше про групповой shift, а не про классическое overfit",
        ],
        8.35,
        5.2,
        3.8,
        1.9,
        size=13,
    )

    # Slide 5: New discovery and next cycle
    s5 = prs.slides.add_slide(blank)
    add_bg(s5, BG)
    add_title(s5, "Новое открытие: color aug вредит", 0.8, 0.45, 7.0, 0.9)
    add_card(s5, 0.8, 1.35, 7.3, 5.95, CARD_2)
    s5.shapes.add_picture(str(chart_ablation), Inches(1.15), Inches(1.8), width=Inches(6.6), height=Inches(3.5))
    delta_acc = ablation["delta_no_color_minus_full"]["acc"]
    delta_f1 = ablation["delta_no_color_minus_full"]["f1_macro"]
    add_subtitle(
        s5,
        f"one-fold A/B (ConvNeXt): no_color лучше на +{delta_acc:.4f} ACC и +{delta_f1:.4f} F1.",
        1.15,
        5.45,
        6.8,
        0.7,
        size=15,
        color=ACCENT,
    )
    add_card(s5, 8.35, 1.35, 4.1, 5.95, CARD)
    add_bullets(
        s5,
        [
            "1) Полный retrain ансамбля в no_color",
            "2) Stage-2 fine-tune (12→20 эпох, сниженный LR)",
            "3) Refit pair-experts + threshold tuning",
            "4) Финальный выбор по public + OOF + group-risk",
            "5) Цель: стабильный top на private split",
        ],
        8.6,
        1.95,
        3.55,
        4.8,
        size=14,
    )

    prs.save(str(OUTPUT_PPTX))
    return OUTPUT_PPTX


def main() -> None:
    if not REPORT_PATH.exists():
        raise FileNotFoundError(f"Missing report markdown: {REPORT_PATH}")
    out = build_presentation()
    print(out)


if __name__ == "__main__":
    main()
