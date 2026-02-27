#!/opt/homebrew/bin/python3.11
from __future__ import annotations

import json
import math
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as patches
import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
ASSETS_DIR = ROOT / "assets"
ARTIFACTS_DIR = ROOT / "artifacts"
DOCS_DIR = ROOT / "docs"
PRES_DIR = DOCS_DIR / "presentation"
PRES_ASSETS_DIR = ASSETS_DIR / "presentation"

OUTPUT_PPTX = PRES_DIR / "DL_Lab1_Grocery_Classification_Final_2026-02-27.pptx"
OUTPUT_NOTES = PRES_DIR / "SPEAKER_NOTES_10MIN_RU.md"


# Palette (warm-neutral + green accent)
BG = RGBColor(248, 246, 241)
TITLE_BAR = RGBColor(28, 45, 64)
TEXT = RGBColor(28, 38, 48)
MUTED = RGBColor(91, 103, 113)
ACCENT = RGBColor(217, 108, 56)
ACCENT_2 = RGBColor(71, 138, 96)
CARD_BG = RGBColor(255, 255, 255)


def emu_to_in(value: int) -> float:
    return value / 914400.0


def in_to_emu(value: float) -> int:
    return int(value * 914400)


@dataclass
class DeckBuilder:
    prs: Presentation

    def __post_init__(self) -> None:
        self.slide_no = 0

    def new_slide(self, title: str, subtitle: str | None = None):
        self.slide_no += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            self.prs.slide_width,
            self.prs.slide_height,
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG
        bg.line.fill.background()

        # Title bar
        bar_h = Inches(0.95)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            0,
            self.prs.slide_width,
            bar_h,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TITLE_BAR
        bar.line.fill.background()

        # Accent strip
        strip = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0,
            bar_h - Inches(0.06),
            self.prs.slide_width,
            Inches(0.06),
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = ACCENT
        strip.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.45),
            Inches(0.16),
            self.prs.slide_width - Inches(1.4),
            Inches(0.48),
        )
        tf = title_box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Calibri"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.47),
                Inches(0.58),
                self.prs.slide_width - Inches(1.4),
                Inches(0.26),
            )
            stf = sub_box.text_frame
            stf.clear()
            sp = stf.paragraphs[0]
            sp.text = subtitle
            sp.font.name = "Calibri"
            sp.font.size = Pt(12)
            sp.font.color.rgb = RGBColor(221, 232, 243)

        # Footer line + slide number
        footer = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.35),
            self.prs.slide_height - Inches(0.32),
            self.prs.slide_width - Inches(0.7),
            Inches(0.01),
        )
        footer.fill.solid()
        footer.fill.fore_color.rgb = RGBColor(214, 214, 214)
        footer.line.fill.background()

        num = slide.shapes.add_textbox(
            self.prs.slide_width - Inches(0.9),
            self.prs.slide_height - Inches(0.28),
            Inches(0.5),
            Inches(0.2),
        )
        ntf = num.text_frame
        ntf.clear()
        np = ntf.paragraphs[0]
        np.text = str(self.slide_no)
        np.alignment = PP_ALIGN.RIGHT
        np.font.name = "Calibri"
        np.font.size = Pt(10)
        np.font.color.rgb = MUTED

        return slide

    def add_notes(self, slide, text: str) -> None:
        notes = slide.notes_slide.notes_text_frame
        notes.clear()
        notes.text = text



def add_card(slide, left: float, top: float, width: float, height: float):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = RGBColor(223, 223, 223)
    card.line.width = Pt(1)
    return card



def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: Iterable[str],
    font_size: int = 21,
    color: RGBColor = TEXT,
    line_space: float = 1.18,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    bullets = list(bullets)
    for idx, txt in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {txt}"
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = 0
        p.line_spacing = line_space
        p.space_after = Pt(8)



def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = TEXT,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color



def add_image_contain(slide, image_path: Path, left: float, top: float, width: float, height: float):
    with Image.open(image_path) as img:
        iw, ih = img.size
    box_w = in_to_emu(width)
    box_h = in_to_emu(height)
    img_ratio = iw / ih
    box_ratio = box_w / box_h

    if img_ratio >= box_ratio:
        out_w = box_w
        out_h = int(out_w / img_ratio)
    else:
        out_h = box_h
        out_w = int(out_h * img_ratio)

    x = in_to_emu(left) + (box_w - out_w) // 2
    y = in_to_emu(top) + (box_h - out_h) // 2
    slide.shapes.add_picture(str(image_path), x, y, width=out_w, height=out_h)

    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(220, 220, 220)
    border.line.width = Pt(1)



def save_timeline_chart(path: Path) -> None:
    milestones = [
        ("2026-02-17", "Baseline CV\n+ first ensembles"),
        ("2026-02-19", "Color A/B\n(no_color win)"),
        ("2026-02-21", "Night zoo\nmodel bakeoff"),
        ("2026-02-22", "Full CV5 hedge\n0.97154 public"),
        ("2026-02-25", "Tinder cleaning\n9889 reviewed"),
        ("2026-02-26", "Mixed old+new\n0.97338 public"),
        ("2026-02-27", "Private results\n0.95200, rank 5/17"),
    ]

    plt.figure(figsize=(12.5, 2.8), dpi=220)
    ax = plt.gca()
    ax.set_xlim(0, len(milestones) - 1)
    ax.set_ylim(-1.2, 1.2)
    ax.axis("off")

    ax.plot([0, len(milestones) - 1], [0, 0], color="#9AA4AE", lw=2)
    for i, (date_s, label) in enumerate(milestones):
        color = "#D96C38" if i in {1, 5, 6} else "#4B8A60"
        ax.scatter(i, 0, s=180, color=color, zorder=3)
        va = "bottom" if i % 2 == 0 else "top"
        y = 0.18 if i % 2 == 0 else -0.2
        ax.text(i, y, date_s, ha="center", va=va, fontsize=9, fontweight="bold", color="#1C2D40")
        ax.text(i, y + (0.27 if i % 2 == 0 else -0.27), label, ha="center", va=va, fontsize=8, color="#334")

    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#F8F6F1")
    plt.close()



def save_public_private_chart(scores_df: pd.DataFrame, path: Path) -> None:
    show_cols = [
        "submission_mixed_lr_mse_geo8_oof_acc.csv",
        "submission_cv5_all20_lr_geo8_equal.csv",
        "submission_mixed_cwr_geo8_oof_acc.csv",
        "attention_ensemble_submission.csv",
        "submission_default_lr_geo8_equal.csv",
    ]
    sdf = scores_df[scores_df["submission_file"].isin(show_cols)].copy()
    order_map = {name: i for i, name in enumerate(show_cols)}
    sdf["_ord"] = sdf["submission_file"].map(order_map)
    sdf = sdf.sort_values("_ord")

    labels = [
        "mixed LR\n(best public)",
        "full CV5\nstable hedge",
        "mixed CWR",
        "attention\n(best private)",
        "partial\nsmoke",
    ]

    x = list(range(len(sdf)))
    public = sdf["public_score"].tolist()
    private = sdf["private_score"].tolist()

    plt.figure(figsize=(12.8, 4.6), dpi=220)
    bw = 0.36
    plt.bar([i - bw / 2 for i in x], public, width=bw, color="#D96C38", label="Public")
    plt.bar([i + bw / 2 for i in x], private, width=bw, color="#4B8A60", label="Private")

    for i, (pv, pr) in enumerate(zip(public, private)):
        plt.text(i - bw / 2, pv + 0.0008, f"{pv:.5f}", ha="center", va="bottom", fontsize=8)
        plt.text(i + bw / 2, pr + 0.0008, f"{pr:.5f}", ha="center", va="bottom", fontsize=8)

    plt.ylim(0.94, 0.978)
    plt.xticks(x, labels, fontsize=9)
    plt.ylabel("Score")
    plt.title("Public vs Private: ключевые сабмиты")
    plt.grid(axis="y", linestyle="--", alpha=0.25)
    plt.legend(loc="upper right")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#FFFFFF")
    plt.close()



def save_gap_chart(scores_df: pd.DataFrame, path: Path) -> None:
    sdf = scores_df.copy()
    sdf["gap"] = sdf["private_score"] - sdf["public_score"]
    sdf = sdf.sort_values("gap")

    labels = []
    for n in sdf["submission_file"]:
        if "mixed_lr_mse_geo8_oof_acc.csv" in n:
            labels.append("mixed LR")
        elif "cv5_all20" in n:
            labels.append("full CV5")
        elif "attention" in n:
            labels.append("attention")
        elif "mixed_cwr" in n:
            labels.append("mixed CWR")
        elif "meta_cwr" in n:
            labels.append("meta CWR")
        elif "default" in n:
            labels.append("partial smoke")
        elif "aggr1" in n:
            labels.append("pair-exp aggr1")
        elif "aggr2" in n:
            labels.append("pair-exp aggr2")
        else:
            labels.append(n[:18])

    colors = ["#C04A2B" if g < 0 else "#2E8B57" for g in sdf["gap"]]

    plt.figure(figsize=(12.8, 4.6), dpi=220)
    y = list(range(len(sdf)))
    plt.barh(y, sdf["gap"], color=colors)
    for yi, gap in zip(y, sdf["gap"]):
        txt = f"{gap:+.5f}"
        x = gap + (0.0007 if gap >= 0 else -0.0007)
        ha = "left" if gap >= 0 else "right"
        plt.text(x, yi, txt, va="center", ha=ha, fontsize=8)

    plt.axvline(0, color="#666", lw=1)
    plt.yticks(y, labels, fontsize=9)
    plt.xlabel("Private - Public")
    plt.title("Public-private сдвиг по сабмитам")
    plt.grid(axis="x", linestyle="--", alpha=0.25)
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#FFFFFF")
    plt.close()



def save_tinder_stacked_chart(class_df: pd.DataFrame, path: Path) -> None:
    cdf = class_df.copy()
    cdf = cdf.sort_values("other_class", ascending=False)

    x = range(len(cdf))
    ok = cdf["ok"].tolist()
    relabel = cdf["other_class"].tolist()
    trash = cdf["trash"].tolist()
    labels = cdf["class_name"].tolist()

    plt.figure(figsize=(13, 5.0), dpi=220)
    plt.bar(x, ok, color="#4B8A60", label="keep")
    plt.bar(x, relabel, bottom=ok, color="#D96C38", label="relabel")
    plt.bar(x, trash, bottom=[a + b for a, b in zip(ok, relabel)], color="#7A8796", label="trash")
    plt.xticks(list(x), labels, rotation=35, ha="right", fontsize=9)
    plt.ylabel("Images")
    plt.title("Tinder-cleaning: действия по классам")
    plt.grid(axis="y", linestyle="--", alpha=0.22)
    plt.legend(loc="upper right")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#FFFFFF")
    plt.close()



def save_tinder_pie(summary_json: dict, path: Path) -> None:
    counts = summary_json["final_decision_counts"]
    labels = ["keep", "relabel", "trash"]
    values = [counts["ok"], counts["other_class"], counts["trash"]]
    colors = ["#4B8A60", "#D96C38", "#7A8796"]

    plt.figure(figsize=(5.2, 4.4), dpi=220)
    plt.pie(
        values,
        labels=labels,
        colors=colors,
        autopct=lambda pct: f"{pct:.1f}%",
        startangle=120,
        textprops={"fontsize": 10},
    )
    plt.title("Итог ручной чистки")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#FFFFFF")
    plt.close()



def save_architecture_diagram(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12.8, 4.8), dpi=220)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 42)
    ax.axis("off")

    def box(x, y, w, h, label, fc, ec="#2C3E50"):
        r = patches.FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle="round,pad=0.02,rounding_size=2.2",
            linewidth=1.2,
            edgecolor=ec,
            facecolor=fc,
        )
        ax.add_patch(r)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=10)

    box(2, 24, 22, 12, "Data Layer\nraw dataset + Tinder cleaning\n(top_new_dataset)", "#EAF4EA")
    box(29, 24, 22, 12, "Training Layer\nConvNeXt-S / EffNetV2-S / DeiT3-S\nno_color + SAM + SWA", "#FCEEDF")
    box(56, 24, 22, 12, "Ensemble Layer\nLR(MSE), geo8 TTA,\nmixed old+new", "#EFE9F8")
    box(83, 24, 15, 12, "Delivery\nKaggle submit\n+ reports", "#E6F0FA")

    box(10, 5, 30, 11, "Repro Layer\nrun scripts + repro/ + docs", "#FFFFFF")
    box(48, 5, 44, 11, "Artifacts Layer\nsubmissions, folds, manual-review logs, weights links", "#FFFFFF")

    arrow = dict(arrowstyle="->", lw=1.8, color="#4B5D70")
    ax.annotate("", xy=(29, 30), xytext=(24, 30), arrowprops=arrow)
    ax.annotate("", xy=(56, 30), xytext=(51, 30), arrowprops=arrow)
    ax.annotate("", xy=(83, 30), xytext=(78, 30), arrowprops=arrow)
    ax.annotate("", xy=(25, 16), xytext=(18, 24), arrowprops=arrow)
    ax.annotate("", xy=(66, 16), xytext=(66, 24), arrowprops=arrow)

    ax.text(2, 39, "Архитектура итогового решения", fontsize=14, weight="bold", color="#1C2D40")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#F8F6F1")
    plt.close()



def save_repro_flow(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(12.8, 3.2), dpi=220)
    ax.set_xlim(0, 100)
    ax.set_ylim(0, 30)
    ax.axis("off")

    steps = [
        (3, "1. Setup env\npip install -r requirements"),
        (24, "2. Check layout\nrepro/scripts/01_check_layout.sh"),
        (45, "3. CV5 hedge\n02_reproduce_full_cv5_submission.sh"),
        (66, "4. Mixed best\n03_reproduce_best_mixed_submission.sh"),
        (87, "5. Verify\nsubmission format + scores"),
    ]

    for i, (x, label) in enumerate(steps):
        r = patches.FancyBboxPatch(
            (x, 8),
            16,
            14,
            boxstyle="round,pad=0.02,rounding_size=2",
            linewidth=1.1,
            edgecolor="#425466",
            facecolor="#FFFFFF",
        )
        ax.add_patch(r)
        ax.text(x + 8, 15, label, ha="center", va="center", fontsize=9)
        if i < len(steps) - 1:
            ax.annotate(
                "",
                xy=(x + 18, 15),
                xytext=(x + 16, 15),
                arrowprops={"arrowstyle": "->", "lw": 1.5, "color": "#5A6B7B"},
            )

    ax.text(2, 26, "Быстрый маршрут воспроизведения", fontsize=13, weight="bold", color="#1C2D40")
    plt.tight_layout()
    plt.savefig(path, bbox_inches="tight", facecolor="#F8F6F1")
    plt.close()



def add_submission_table(slide, scores_df: pd.DataFrame) -> None:
    view = scores_df.copy()
    view = view[["submission_file", "public_score", "private_score", "category"]]
    view = view.sort_values("public_score", ascending=False)
    view = view.head(8)

    rows = len(view) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.6), Inches(1.3), Inches(12.2), Inches(5.4)).table

    table.columns[0].width = Inches(4.9)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    table.columns[3].width = Inches(3.7)

    headers = ["Submission", "Public", "Private", "Category"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TITLE_BAR
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)

    for r, (_, row) in enumerate(view.iterrows(), start=1):
        values = [
            str(row["submission_file"]),
            f"{row['public_score']:.5f}",
            f"{row['private_score']:.5f}",
            str(row["category"]),
        ]
        for c, val in enumerate(values):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT
            if c in (1, 2):
                p.alignment = PP_ALIGN.CENTER
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(246, 246, 246)



def build_deck() -> None:
    PRES_DIR.mkdir(parents=True, exist_ok=True)
    PRES_ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    scores_df = pd.read_csv(ARTIFACTS_DIR / "submissions" / "public_scores_registry.csv")
    class_df = pd.read_csv(ARTIFACTS_DIR / "manual_review" / "per_class_stats.csv")
    summary_json = json.loads((ARTIFACTS_DIR / "manual_review" / "tinder_session_export_summary.json").read_text(encoding="utf-8"))

    # Generated charts
    timeline_png = PRES_ASSETS_DIR / "timeline_milestones.png"
    pp_png = PRES_ASSETS_DIR / "public_private_key_submissions.png"
    gap_png = PRES_ASSETS_DIR / "public_private_gap.png"
    class_png = PRES_ASSETS_DIR / "tinder_actions_by_class.png"
    pie_png = PRES_ASSETS_DIR / "tinder_actions_pie.png"
    arch_png = PRES_ASSETS_DIR / "solution_architecture.png"
    repro_png = PRES_ASSETS_DIR / "repro_flow.png"

    save_timeline_chart(timeline_png)
    save_public_private_chart(scores_df, pp_png)
    save_gap_chart(scores_df, gap_png)
    save_tinder_stacked_chart(class_df, class_png)
    save_tinder_pie(summary_json, pie_png)
    save_architecture_diagram(arch_png)
    save_repro_flow(repro_png)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    b = DeckBuilder(prs)

    notes_lines = []

    # 1
    slide = b.new_slide(
        "DL Lab 1 Grocery Classification: итоговое решение",
        "Команда: Фёдор Грачёв, Ярослав Кулизнев, Константин Родионов",
    )
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_text(
        slide,
        0.9,
        1.7,
        11.5,
        1.0,
        "Финальные результаты: public 0.97338 | private 0.95200 (5/17)",
        size=28,
        bold=True,
        color=TITLE_BAR,
    )
    add_text(
        slide,
        0.95,
        2.85,
        11.3,
        2.8,
        "Период: февраль 2026\nСоревнование: dl-lab-1-image-classification (15 классов фруктов/овощей)\n\nВ презентации: путь от базового пайплайна до финального mixed old+new решения,\nчто реально сработало, что не сработало, и как это воспроизвести из репозитория.",
        size=18,
        color=TEXT,
    )
    b.add_notes(slide, "20 сек. Открываю доклад: задача, команда, финальные метрики и структура выступления.")
    notes_lines.append("1. 00:00-00:20 — Титул, финальные метрики и план.")

    # 2
    slide = b.new_slide("План доклада (~10 минут)")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.9,
        1.55,
        11.4,
        5.0,
        [
            "Контекст задачи, ограничения и роли в команде",
            "Эволюция решения: baseline -> CV5 hedge -> mixed old+new",
            "Ключевые эксперименты: что улучшило и что ухудшило",
            "Tinder-style data cleaning и подтверждение эффекта",
            "Public/private сдвиг и финальный private результат",
            "Репозиторий и пошаговая воспроизводимость",
            "Идеи на post-deadline улучшения",
        ],
        font_size=20,
    )
    b.add_notes(slide, "20 сек. Проговариваю маршрут, чтобы преподавателю было понятно, что будет системный разбор.")
    notes_lines.append("2. 00:20-00:40 — План выступления.")

    # 3
    slide = b.new_slide("Постановка задачи и ограничения")
    add_card(slide, 0.6, 1.2, 6.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        5.5,
        5.0,
        [
            "15 классов, target: точная классификация grocery-изображений",
            "Основная метрика соревнования: leaderboard score",
            "Жёсткий дедлайн и распределённое обучение на MPS/CUDA/Colab",
            "Тестовый набор не меняли, все изменения только в train",
            "Дополнительный риск: public/private domain shift",
        ],
        font_size=18,
    )
    add_card(slide, 7.0, 1.2, 5.7, 5.6)
    add_text(
        slide,
        7.35,
        1.6,
        5.0,
        4.8,
        "Рабочая цель:\nполучить максимум качества без потери воспроизводимости.\n\nПоэтому финальное решение выбирали не по одному числу public LB,\nа по совокупности:\n• OOF-метрики\n• controlled A/B\n• устойчивость к сдвигу\n• инженерная повторяемость пайплайна.",
        size=18,
    )
    b.add_notes(slide, "20 сек. Здесь важно показать зрелый инженерный подход, а не охоту только за public score.")
    notes_lines.append("3. 00:40-01:00 — Задача, ограничения, критерии выбора.")

    # 4
    slide = b.new_slide("Командный вклад")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.9,
        1.55,
        11.5,
        5.0,
        [
            "Фёдор Грачёв: основной экспериментальный и интеграционный контур, ключевые тренировочные циклы, финальная сборка, репозиторий и презентационные материалы",
            "Ярослав Кулизнев: активная генерация и обсуждение гипотез, проверка альтернативных направлений улучшения пайплайна",
            "Константин Родионов: практическая работа с данными и участие в запусках/сопровождении обучения",
            "Роли пересекались, но финальный релиз собран в едином воспроизводимом контуре",
        ],
        font_size=18,
    )
    b.add_notes(slide, "20 сек. Мягко и профессионально фиксирую роли без конфронтации.")
    notes_lines.append("4. 01:00-01:20 — Распределение ролей в команде.")

    # 5
    slide = b.new_slide("Хронология проекта")
    add_card(slide, 0.55, 1.3, 12.2, 5.4)
    add_image_contain(slide, timeline_png, 0.8, 1.8, 11.7, 2.6)
    add_bullets(
        slide,
        0.95,
        4.7,
        11.4,
        1.7,
        [
            "От baseline-ансамбля -> к full CV5 hedge -> к data-cleaning -> к mixed old+new финалу",
            "После публикации private: подтверждён итог 0.95200 и rank 5/17",
        ],
        font_size=16,
    )
    b.add_notes(slide, "25 сек. Кратко показываю последовательность решений и контрольных точек.")
    notes_lines.append("5. 01:20-01:45 — Временная шкала ключевых этапов.")

    # 6
    slide = b.new_slide("Архитектура итогового решения")
    add_card(slide, 0.55, 1.25, 12.2, 5.45)
    add_image_contain(slide, arch_png, 0.8, 1.6, 11.7, 4.7)
    b.add_notes(slide, "25 сек. Объясняю архитектуру слоями: данные, обучение, ансамбль, доставка, воспроизводимость.")
    notes_lines.append("6. 01:45-02:10 — Архитектура решения (high-level).")

    # 7
    slide = b.new_slide("Эволюция ансамбля: от baseline к сильному стеку")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, ASSETS_DIR / "chart_ensemble_progress.png", 0.85, 1.55, 7.8, 4.9)
    add_bullets(
        slide,
        8.95,
        1.75,
        3.45,
        4.8,
        [
            "Быстрые приросты: weight search + class-bias",
            "Pair-experts дал локальный плюс",
            "Meta-ветки усилили OOF",
            "Финал: late-stage mixed old+new",
        ],
        font_size=15,
    )
    b.add_notes(slide, "20 сек. Поясняю, что рост был ступенчатым и проверяемым, а не случайным.")
    notes_lines.append("7. 02:10-02:30 — График эволюции ансамбля.")

    # 8
    slide = b.new_slide("Ключевой A/B: color-aug vs no_color")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, ASSETS_DIR / "chart_color_ablation.png", 0.85, 1.5, 7.8, 5.0)
    add_bullets(
        slide,
        8.9,
        1.7,
        3.5,
        4.9,
        [
            "ΔACC: +0.00917",
            "ΔF1-macro: +0.00731",
            "no_color стал базой для следующих циклов",
            "Вывод: агрессивный color-jitter вредил классовой семантике",
        ],
        font_size=15,
    )
    b.add_notes(slide, "25 сек. Это самый важный controlled A/B, после него поменялась вся ветка обучения.")
    notes_lines.append("8. 02:30-02:55 — A/B no_color, ключевой перелом.")

    # 9
    slide = b.new_slide("Phase1 probes: что закрепили в дедлайн-конфиге")
    add_card(slide, 0.6, 1.25, 6.1, 5.45)
    add_bullets(
        slide,
        0.95,
        1.6,
        5.35,
        5.0,
        [
            "SWA on, ранний старт swa_start_epoch=5",
            "SAM on",
            "label_smoothing=0.03",
            "weighted sampler",
            "mixup + cutmix",
            "no_color augmentation policy",
        ],
        font_size=18,
    )
    add_card(slide, 6.95, 1.25, 5.75, 5.45)
    add_text(
        slide,
        7.3,
        1.65,
        5.0,
        4.9,
        "One-fold probes подтвердили, что комбинация SAM+SWA+LS=0.03 даёт лучший устойчивый баланс.\n\nТакже отсекли ветки, где отключение SAM/SWA ухудшало метрики.\n\nИтог: дедлайн-конфиг опирался только на подтверждённые изменения.",
        size=17,
    )
    b.add_notes(slide, "20 сек. Проговариваю, что финальный конфиг не случайный — он отфильтрован пробами.")
    notes_lines.append("9. 02:55-03:15 — Принятые гиперпараметры из probes.")

    # 10
    slide = b.new_slide("Full CV5 stable hedge")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.9,
        1.55,
        11.4,
        5.0,
        [
            "20 запусков: 4 модели x 5 фолдов",
            "Внутри фолда: LinearRegression(MSE), positive convex weights",
            "TTA: geo8, fold aggregation: equal",
            "Стабильный результат: submission_cv5_all20_lr_geo8_equal.csv -> 0.97154 public",
            "Этот хедж использовали как контрольную устойчивую точку",
        ],
        font_size=20,
    )
    b.add_notes(slide, "20 сек. Здесь показываю базовый сильный и воспроизводимый anchor-результат.")
    notes_lines.append("10. 03:15-03:35 — Full CV5 как стабильный baseline-хедж.")

    # 11
    slide = b.new_slide("Late-stage mixed old+new: идея и реализация")
    add_card(slide, 0.6, 1.2, 6.1, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        5.35,
        5.1,
        [
            "Смешение pre-clean old zoo и post-clean new zoo",
            "Выравнивание OOF по image_id",
            "Blending: mixed_lr_mse и mixed_cwr",
            "Генерация через mixed_old_new_orchestrator_submit.py",
            "Ключевой выигрыш пришёл из diversity двух пулов",
        ],
        font_size=16,
    )
    add_card(slide, 6.95, 1.2, 5.75, 5.6)
    add_text(
        slide,
        7.3,
        1.65,
        5.0,
        4.9,
        "Best deadline submission:\nsubmission_mixed_lr_mse_geo8_oof_acc.csv\n\nPublic: 0.97338\nPrivate: 0.94675\n\nLR(MSE) оказался надёжнее CWR в итоговом публичном отборе.",
        size=18,
        bold=False,
    )
    b.add_notes(slide, "20 сек. Поясняю механику mixed-ветки и почему она дала максимум на public.")
    notes_lines.append("11. 03:35-03:55 — mixed old+new orchestrator.")

    # 12
    slide = b.new_slide("Public vs Private по ключевым сабмитам")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, pp_png, 0.8, 1.55, 11.7, 4.9)
    b.add_notes(slide, "25 сек. На графике видно, что лучший public не всегда лучший private; внимание на attention-сабмит.")
    notes_lines.append("12. 03:55-04:20 — Сопоставление public/private результатов.")

    # 13
    slide = b.new_slide("Анализ public-private сдвига")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, gap_png, 0.8, 1.55, 8.2, 4.9)
    add_bullets(
        slide,
        9.2,
        1.8,
        3.2,
        4.7,
        [
            "Сильный negative shift у mixed LR",
            "Лучший private у attention ensemble",
            "Итоговая стратегия: держать стабильный CV5-хедж как противовес",
            "Private фиксирует риск переоптимизации под public",
        ],
        font_size=15,
    )
    b.add_notes(slide, "25 сек. Прямо говорю о переоптимизации под public и как это учитывали в выводах.")
    notes_lines.append("13. 04:20-04:45 — Gap-анализ по сабмитам.")

    # 14
    slide = b.new_slide("Официальный private итог")
    add_card(slide, 0.55, 1.25, 12.2, 5.45)
    add_text(
        slide,
        0.95,
        1.75,
        11.4,
        1.0,
        "Команда «батчсайз не влез»: private score 0.95200, rank 5/17",
        size=27,
        bold=True,
        color=TITLE_BAR,
    )
    add_bullets(
        slide,
        1.0,
        3.05,
        11.3,
        3.3,
        [
            "Данные зафиксированы после публикации private leaderboard (2026-02-27)",
            "Финальный репозиторий дополнен post-competition анализом и таблицами",
            "Сильная сторона решения: воспроизводимость всех ключевых этапов",
        ],
        font_size=20,
    )
    b.add_notes(slide, "20 сек. Коротко фиксирую официальный результат и его место в финальном отчёте.")
    notes_lines.append("14. 04:45-05:05 — Private leaderboard итог.")

    # 15
    slide = b.new_slide("Risk profile: почему смотрели не только на public")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, ASSETS_DIR / "chart_risk_profile.png", 0.85, 1.55, 7.8, 4.9)
    add_bullets(
        slide,
        8.95,
        1.85,
        3.4,
        4.7,
        [
            "На IID holdout всё выглядит лучше, чем на group-holdout",
            "Риск концентрируется в некоторых группах (`plu`)",
            "Это объясняет часть private-просадки у агрессивных late-stage стратегий",
        ],
        font_size=15,
    )
    b.add_notes(slide, "20 сек. Связываю private-shift с групповой структурой данных.")
    notes_lines.append("15. 05:05-05:25 — Риск-профиль и group-shift.")

    # 16
    slide = b.new_slide("Tinder-style cleaning: зачем и как")
    add_card(slide, 0.6, 1.2, 6.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        5.4,
        5.0,
        [
            "Локальный web UI с быстрыми решениями: keep / relabel / trash",
            "Есть undo и экспорт действий в replay-совместимый формат",
            "Авто-применение действий к top_new_dataset",
            "Пост-валидация целостности датасета",
            "Полные логи и CSV-артефакты сохранены в репозитории",
        ],
        font_size=16,
    )
    add_card(slide, 6.95, 1.2, 5.75, 5.6)
    add_image_contain(slide, repro_png, 7.2, 1.65, 5.2, 4.6)
    b.add_notes(slide, "20 сек. Объясняю, что это не ручная магия, а воспроизводимый data-engineering pipeline.")
    notes_lines.append("16. 05:25-05:45 — Зачем Tinder-cleaning и как он устроен.")

    # 17
    slide = b.new_slide("Tinder UI: рабочий интерфейс (live screenshot)")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_image_contain(slide, ASSETS_DIR / "tinder_ui_main.png", 0.8, 1.5, 11.7, 5.0)
    b.add_notes(slide, "20 сек. Показываю реальный интерфейс из рабочего скрина, что review делался инструментально.")
    notes_lines.append("17. 05:45-06:05 — Скриншот Tinder UI (основной экран).")

    # 18
    slide = b.new_slide("Tinder UI: модальное решение и работа по классам")
    add_card(slide, 0.55, 1.2, 5.95, 5.6)
    add_card(slide, 6.82, 1.2, 5.95, 5.6)
    add_image_contain(slide, ASSETS_DIR / "tinder_ui_modal.png", 0.78, 1.55, 5.5, 4.7)
    add_image_contain(slide, ASSETS_DIR / "tinder_ui_class_switch.png", 7.05, 1.55, 5.5, 4.7)
    add_text(slide, 0.95, 6.25, 5.2, 0.35, "Modal: trash / relabel", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 7.2, 6.25, 5.2, 0.35, "Class switch + progress", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    b.add_notes(slide, "20 сек. Фиксирую, что интерфейс поддерживает быстрые решения и контроль прогресса по классам.")
    notes_lines.append("18. 06:05-06:25 — Скрины модального действия и class-switch.")

    # 19
    slide = b.new_slide("Результат ручной чистки: количественный срез")
    add_card(slide, 0.55, 1.2, 4.0, 5.6)
    add_image_contain(slide, pie_png, 0.75, 1.6, 3.6, 2.8)
    add_bullets(
        slide,
        0.85,
        4.5,
        3.5,
        2.1,
        [
            "Всего reviewed: 9889",
            "keep: 9696",
            "relabel: 139",
            "trash: 54",
        ],
        font_size=14,
    )
    add_card(slide, 4.8, 1.2, 7.95, 5.6)
    add_image_contain(slide, class_png, 5.1, 1.55, 7.35, 5.0)
    b.add_notes(slide, "25 сек. На цифрах показываю объём и распределение ручной чистки.")
    notes_lines.append("19. 06:25-06:50 — Количественные результаты Tinder-cleaning.")

    # 20
    slide = b.new_slide("Проверка целостности после cleaning")
    add_card(slide, 0.55, 1.25, 12.2, 5.45)
    add_bullets(
        slide,
        0.95,
        1.7,
        11.3,
        4.9,
        [
            "top_new_dataset/train.csv: 9835 строк",
            "Активных изображений на диске: 9835",
            "1:1 соответствие CSV и файловой структуры",
            "mismatch label vs folder: 0",
            "duplicate image_id: 0",
            "Все действия сохранились в artifacts/manual_review/*",
        ],
        font_size=20,
    )
    b.add_notes(slide, "20 сек. Подчёркиваю, что очистка данных была проверена на структурную корректность.")
    notes_lines.append("20. 06:50-07:10 — Data integrity после cleaning.")

    # 21
    slide = b.new_slide("Что ухудшало или не дало прироста")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        11.4,
        5.0,
        [
            "Сильные color-aug: подтверждённая деградация по ACC/F1",
            "Stage4 confidence-aware finetune (текущая реализация): прирост не подтверждён",
            "mixed CWR как финальный public-сабмит: хуже mixed LR",
            "Pair-experts поверх лучшего mixed LR: не добавил прироста к 0.97338",
            "Partial final3 default ensemble: smoke-результат, не боевой",
        ],
        font_size=19,
    )
    b.add_notes(slide, "25 сек. Важный слайд по требованиям: честно фиксирую неудачные ветки и почему они были отбракованы.")
    notes_lines.append("21. 07:10-07:35 — Что не сработало и почему.")

    # 22
    slide = b.new_slide("Воспроизводимость: что есть в репозитории")
    add_card(slide, 0.6, 1.2, 6.1, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        5.3,
        5.0,
        [
            "docs/EXPERIMENTS_RU.md — полный экспериментальный срез",
            "docs/REPRODUCIBILITY_RU.md — сценарии воспроизведения",
            "repro/README.md + 3 скрипта быстрого маршрута",
            "artifacts/submissions/* — финальные CSV и реестр оценок",
            "artifacts/manual_review/* — все логи Tinder-cleaning",
        ],
        font_size=16,
    )
    add_card(slide, 6.95, 1.2, 5.75, 5.6)
    add_text(
        slide,
        7.25,
        1.65,
        5.1,
        4.9,
        "Быстрый путь:\n1) pip install -r requirements_common_team_train.txt\n2) bash repro/scripts/01_check_layout.sh\n3) bash repro/scripts/02_reproduce_full_cv5_submission.sh\n4) bash repro/scripts/03_reproduce_best_mixed_submission.sh\n\nВеса и бандлы вынесены в Google Drive, ссылки и SHA256 в weights/README.md.",
        size=16,
    )
    b.add_notes(slide, "20 сек. Показываю, что преподаватель может повторить пайплайн пошагово.")
    notes_lines.append("22. 07:35-07:55 — Репозиторий и reproducibility.")

    # 23
    slide = b.new_slide("Внешние артефакты (Google Drive)")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_text(
        slide,
        0.9,
        1.55,
        11.4,
        4.9,
        "Опубликованы архивы для воспроизведения:\n\n"
        "• team_bundle_top_new_tinder_plus_ckpts_2026-02-26_v1.zip (4.74 GB)\n"
        "• team_final3_cv5_split_bundle_2026-02-26_v1.zip (724.57 MB)\n"
        "• convnext_small_deadline14f4_weights_f0f1_2026-02-26.zip (755.31 MB)\n\n"
        "Папка: https://drive.google.com/open?id=1Jgb1xTOhgvdicsL0oPyxplWl9yIP98-b\n"
        "Все ссылки и SHA256 доступны в weights/README.md и repro/upload_manifest_google_drive.json",
        size=17,
    )
    b.add_notes(slide, "18 сек. Коротко фиксирую, где лежат большие артефакты и как проверяется целостность.")
    notes_lines.append("23. 07:55-08:13 — Веса и внешние бандлы.")

    # 24
    slide = b.new_slide("Идеи улучшения после дедлайна")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        11.4,
        5.0,
        [
            "Расширить diversity для mixed LR (новые независимые источники вероятностей)",
            "Полноценно перепроверить CWR/attention на полном покрытии фолдов",
            "CatBoost meta на честном OOF без утечек",
            "Group-aware отбор финального сабмита (по worst-group риску)",
            "High-resolution stage-2 fine-tune (320/384) для fine-grained классов",
            "Refit pair-experts в новом data regime",
        ],
        font_size=18,
    )
    b.add_notes(slide, "20 сек. Показываю, что дорожная карта улучшений есть и она реалистична.")
    notes_lines.append("24. 08:13-08:33 — Post-deadline roadmap.")

    # 25
    slide = b.new_slide("Итоги")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.75,
        11.3,
        4.8,
        [
            "Сильный результат достигнут комбинацией data cleaning + устойчивого тренинга + ансамблирования",
            "Ключевые решения подтверждены controlled A/B и Kaggle-проверками",
            "Финальный репозиторий и артефакты готовы для воспроизведения и проверки преподавателем",
            "Официальный private итог: 0.95200 (rank 5/17)",
        ],
        font_size=20,
    )
    add_text(
        slide,
        0.95,
        6.15,
        11.3,
        0.3,
        "Репозиторий: https://github.com/HVGrach/dl-lab1-grocery-classification-solution",
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    b.add_notes(slide, "25 сек. Финальный слайд: результат, воспроизводимость, и готовность к проверке.")
    notes_lines.append("25. 08:33-08:58 — Заключение.")

    # 26 backup
    slide = b.new_slide("Backup: таблица ключевых сабмитов")
    add_submission_table(slide, scores_df)
    b.add_notes(slide, "Backup. По запросу можно пройти по каждой строке: файл, category, public/private.")

    # 27 backup
    slide = b.new_slide("Backup: карта основных скриптов")
    add_card(slide, 0.55, 1.2, 12.2, 5.6)
    add_bullets(
        slide,
        0.95,
        1.55,
        11.4,
        5.0,
        [
            "dl_lab1/scripts/mixed_old_new_orchestrator_submit.py — late-stage mixed сабмиты",
            "dl_lab1/scripts/top_new_final3_cv5_train_orchestrator.py — split-training дедлайн-контура",
            "dl_lab1/scripts/adaptive_top_new_night_pipeline.py — phase1 probes",
            "dl_lab1/scripts/dataset_tinder_review_app.py — интерфейс ручной чистки",
            "repro/scripts/0*.sh — быстрый маршрут воспроизведения",
        ],
        font_size=18,
    )
    b.add_notes(slide, "Backup. Карта скриптов нужна, если преподаватель спросит «что запускать и где».")

    # 28 backup
    slide = b.new_slide("Backup: reference screenshots")
    add_card(slide, 0.55, 1.2, 5.95, 5.6)
    add_card(slide, 6.82, 1.2, 5.95, 5.6)
    add_image_contain(slide, ASSETS_DIR / "cleaning_preview.jpg", 0.75, 1.55, 5.55, 4.8)
    add_image_contain(slide, ASSETS_DIR / "aug_debug_cover.png", 7.02, 1.55, 5.55, 4.8)
    add_text(slide, 1.0, 6.2, 5.0, 0.25, "cleaning preview", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, 7.2, 6.2, 5.0, 0.25, "augmentation debug cover", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    b.add_notes(slide, "Backup. Дополнительные визуалы по data-cleaning и augmentation.")

    prs.save(OUTPUT_PPTX)

    # Speaker notes markdown
    total_time_sec = 0
    for line in notes_lines:
        # parse mm:ss range from line text
        try:
            rng = line.split("—")[0].split()[-1]
            # no strict parsing; keep manually assigned below
        except Exception:
            pass

    notes_md = [
        "# Speaker Notes (~10 minutes)",
        "",
        f"Презентация: `{OUTPUT_PPTX.name}`",
        "",
        "Ниже базовый тайминг по основным 25 слайдам (backup не включены).",
        "",
    ]
    notes_md.extend(notes_lines)
    notes_md.append("")
    notes_md.append("Итоговый целевой хронометраж: ~9 минут.")
    notes_md.append("Если нужно ровно 10 минут: добавить 1 минуту на вопросы по backup-слайдам 26-28.")

    OUTPUT_NOTES.write_text("\n".join(notes_md), encoding="utf-8")

    print(f"Saved presentation: {OUTPUT_PPTX}")
    print(f"Saved notes: {OUTPUT_NOTES}")
    print(f"Generated chart assets in: {PRES_ASSETS_DIR}")


if __name__ == "__main__":
    build_deck()
